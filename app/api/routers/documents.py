"""
文書生成 APIルーター
Word / Excel / PowerPoint ファイルを生成して返す
"""
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from loguru import logger
import tempfile
import os
from pathlib import Path

from services.rag_engine import rag_engine

router = APIRouter()


class DocumentRequest(BaseModel):
    doc_type: str       # "excel" / "word" / "pptx"
    requirements: str   # 作成要件（自然言語）
    search_query: str   # 参照資料の検索キーワード
    template: str = "default"


@router.post("/generate")
async def generate_document(request: DocumentRequest):
    """文書ドラフトを生成してダウンロードさせる"""
    logger.info(f"文書生成リクエスト: type={request.doc_type}")

    # まずRAGでドラフトコンテンツを生成
    rag_result = await rag_engine.generate_document_draft(
        requirements=request.requirements,
        doc_type=request.doc_type,
        query=request.search_query,
    )

    draft_content = rag_result.answer

    try:
        if request.doc_type == "word":
            file_path = _generate_word(draft_content, request.requirements)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename = "draft_report.docx"

        elif request.doc_type == "excel":
            file_path = _generate_excel(draft_content, request.requirements)
            media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            filename = "draft_sheet.xlsx"

        elif request.doc_type == "pptx":
            file_path = _generate_pptx(draft_content, request.requirements)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = "draft_presentation.pptx"

        else:
            raise HTTPException(status_code=400, detail=f"未対応のdoc_type: {request.doc_type}")

        return FileResponse(
            path=file_path,
            media_type=media_type,
            filename=filename,
        )
    except Exception as e:
        logger.error(f"文書生成エラー: {e}")
        raise HTTPException(status_code=500, detail=str(e))


def _generate_word(content: str, title: str) -> str:
    """Word文書を生成してパスを返す"""
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # タイトル
    heading = doc.add_heading(title[:50], level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # 本文（マークダウンを簡易変換）
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            doc.add_paragraph()
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("- ") or line.startswith("* "):
            p = doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)

    # 一時ファイルに保存
    tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
    doc.save(tmp.name)
    return tmp.name


def _generate_excel(content: str, title: str) -> str:
    """見積書フォーマットのExcelファイルを生成してパスを返す

    content には ESTIMATE_EXTRACTION_PROMPT で抽出させたJSON文字列が渡される想定。
    JSONが見つからない/不正な場合は空の見積書テンプレートにフォールバックする。
    """
    import json
    import re
    from datetime import datetime, timedelta
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    # --- JSONパース（LLMの出力からJSONを抽出）---
    json_match = re.search(r'\{.*\}', content, re.DOTALL)
    if json_match:
        try:
            data = json.loads(json_match.group())
        except json.JSONDecodeError:
            data = {}
    else:
        data = {}

    # デフォルト値のフォールバック
    today = datetime.today()
    to_company    = data.get("to_company") or "〇〇株式会社 御中"
    to_address    = data.get("to_address") or ""
    estimate_date = data.get("estimate_date") or today.strftime("%Y-%m-%d")
    estimate_no   = data.get("estimate_no") or "001"
    valid_until   = data.get("valid_until") or (today + timedelta(days=30)).strftime("%Y-%m-%d")
    delivery      = data.get("delivery_date") or "要相談"
    payment       = data.get("payment_terms") or "月末締め翌月末払い"
    from_company  = data.get("from_company") or "●●株式会社"
    from_address  = data.get("from_address") or ""
    from_tel      = data.get("from_tel") or ""
    from_email    = data.get("from_email") or ""
    items         = data.get("items") or []
    notes         = data.get("notes") or ""

    wb = Workbook()
    ws = wb.active
    ws.title = "見積書"

    # --- 列幅設定（A〜R列、サンプル準拠）---
    col_widths = {
        "A": 3, "B": 18, "C": 4, "D": 4, "E": 4,
        "F": 4, "G": 4, "H": 4, "I": 4, "J": 6,
        "K": 6, "L": 6, "M": 5, "N": 10, "O": 4,
        "P": 6, "Q": 14, "R": 3,
    }
    for col, width in col_widths.items():
        ws.column_dimensions[col].width = width

    # --- スタイル定義 ---
    label_font  = Font(name="メイリオ", size=10)
    bold_font   = Font(name="メイリオ", bold=True, size=10)
    normal_font = Font(name="メイリオ", size=10)
    header_fill = PatternFill("solid", fgColor="1F4E79")  # 紺
    center = Alignment(horizontal="center", vertical="center")
    right = Alignment(horizontal="right", vertical="center")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    # --- タイトル行 (row=1) ---
    ws.merge_cells("B1:R1")
    ws["B1"] = "見　積　書"
    ws["B1"].font = Font(name="メイリオ", bold=True, size=20)
    ws["B1"].alignment = center
    ws.row_dimensions[1].height = 36

    # --- 見積日・見積番号 (row=2,3) ---
    ws["N2"] = "見積日"; ws["N2"].font = bold_font
    ws["P2"] = estimate_date; ws["P2"].font = normal_font
    ws["N3"] = "見積番号"; ws["N3"].font = bold_font
    ws["P3"] = estimate_no; ws["P3"].font = normal_font

    # --- 宛先 (row=4〜7) ---
    ws["B4"] = to_company
    ws["B4"].font = Font(name="メイリオ", bold=True, size=13)
    ws.merge_cells("B4:M4")
    ws["B5"] = to_address
    ws["B5"].font = normal_font
    ws["B7"] = "下記の通り、お見積り申し上げます。"
    ws["B7"].font = normal_font

    # --- 発行者情報（右側 row=6〜9）---
    ws["N6"] = from_company; ws["N6"].font = bold_font
    ws.merge_cells("N6:R6")
    ws["N7"] = from_address; ws["N7"].font = normal_font
    ws.merge_cells("N7:R7")
    ws["N8"] = f"電話：{from_tel}"; ws["N8"].font = normal_font
    ws["N9"] = f"メール：{from_email}"; ws["N9"].font = normal_font

    # --- 合計金額ブロック (row=8〜12) ---
    ws["B8"] = "お見積り金額（税込）"; ws["B8"].font = bold_font
    ws.merge_cells("B8:G8")
    ws["B10"] = "納期"; ws["B10"].font = label_font
    ws["C10"] = delivery; ws["C10"].font = normal_font
    ws["B11"] = "支払条件"; ws["B11"].font = label_font
    ws["C11"] = payment; ws["C11"].font = normal_font
    ws["B12"] = "有効期限"; ws["B12"].font = label_font
    ws["C12"] = valid_until; ws["C12"].font = normal_font

    # --- 明細ヘッダー (row=13) ---
    header_row = 13
    headers = [
        ("B", "内容"), ("K", "軽減"), ("L", "数量"),
        ("M", "単位"), ("N", "単価（税抜）"), ("P", "税率"),
        ("Q", "金額（税抜）"),
    ]
    for col, label in headers:
        cell = ws[f"{col}{header_row}"]
        cell.value = label
        cell.font = Font(name="メイリオ", bold=True, size=9, color="FFFFFF")
        cell.fill = header_fill
        cell.alignment = center
        cell.border = thin_border
    ws.merge_cells(f"B{header_row}:J{header_row}")
    ws.row_dimensions[header_row].height = 18

    # --- 明細行（row=14〜） ---
    item_start = 14
    item_rows = max(len(items), 10)  # 最低10行確保
    for i in range(item_rows):
        r = item_start + i
        ws.row_dimensions[r].height = 18
        ws.merge_cells(f"B{r}:J{r}")
        if i < len(items):
            item = items[i]
            qty = item.get("qty", 0)
            unit_price = item.get("unit_price", 0)
            tax_rate = item.get("tax_rate", 0.10)

            ws[f"B{r}"] = item.get("name", ""); ws[f"B{r}"].font = normal_font
            ws[f"K{r}"] = "※" if tax_rate == 0.08 else ""
            ws[f"K{r}"].alignment = center
            ws[f"L{r}"] = qty
            ws[f"L{r}"].alignment = center
            ws[f"M{r}"] = item.get("unit", "式")
            ws[f"M{r}"].alignment = center
            ws[f"N{r}"] = unit_price
            ws[f"N{r}"].number_format = "#,##0"
            ws[f"N{r}"].alignment = right
            ws[f"P{r}"] = tax_rate
            ws[f"P{r}"].number_format = "0%"
            ws[f"P{r}"].alignment = center
            ws[f"Q{r}"] = f'=IF(N{r}="","",L{r}*N{r})'
            ws[f"Q{r}"].number_format = "#,##0"
        else:
            # 空行
            ws[f"Q{r}"] = f'=IF(N{r}="","",L{r}*N{r})'
            ws[f"Q{r}"].number_format = "#,##0"

        # 罫線
        for col in ["B", "K", "L", "M", "N", "O", "P", "Q"]:
            ws[f"{col}{r}"].border = thin_border

    # --- 小計・税額・合計 ---
    sum_row = item_start + item_rows
    tax8_row = sum_row + 1
    tax10_row = sum_row + 2
    total_row = sum_row + 3

    item_range_p = f"P{item_start}:P{item_start + item_rows - 1}"
    item_range_q = f"Q{item_start}:Q{item_start + item_rows - 1}"

    ws[f"N{sum_row}"] = "小計（税抜）"; ws[f"N{sum_row}"].font = bold_font
    ws[f"Q{sum_row}"] = f"=SUM({item_range_q})"
    ws[f"Q{sum_row}"].number_format = "#,##0"

    # N()でラップし、空行のQが文字列""でも#VALUE!エラーにならないようにする
    ws[f"N{tax8_row}"] = "消費税（8%）"; ws[f"N{tax8_row}"].font = bold_font
    ws[f"Q{tax8_row}"] = (
        f"=ROUND(SUMPRODUCT(({item_range_p}=0.08)*N({item_range_q}))*0.08,0)"
    )
    ws[f"Q{tax8_row}"].number_format = "#,##0"

    ws[f"N{tax10_row}"] = "消費税（10%）"; ws[f"N{tax10_row}"].font = bold_font
    ws[f"Q{tax10_row}"] = (
        f"=ROUND(SUMPRODUCT(({item_range_p}=0.10)*N({item_range_q}))*0.10,0)"
    )
    ws[f"Q{tax10_row}"].number_format = "#,##0"

    ws[f"N{total_row}"] = "合計（税込）"
    ws[f"N{total_row}"].font = bold_font
    ws[f"Q{total_row}"] = f"=Q{sum_row}+Q{tax8_row}+Q{tax10_row}"
    ws[f"Q{total_row}"].number_format = "#,##0"
    ws[f"Q{total_row}"].font = bold_font

    # 合計金額を宛先下の大きいセルに参照
    ws["B9"] = f"=Q{total_row}"
    ws["B9"].number_format = "¥#,##0"
    ws["B9"].font = Font(name="メイリオ", bold=True, size=14)
    ws.merge_cells("B9:G9")

    # 備考
    if notes:
        note_row = total_row + 2
        ws[f"B{note_row}"] = "備考"
        ws[f"B{note_row}"].font = bold_font
        ws[f"C{note_row}"] = notes
        ws[f"C{note_row}"].font = normal_font

    tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
    wb.save(tmp.name)
    return tmp.name


def _generate_pptx(content: str, title: str) -> str:
    """PowerPoint文書を生成してパスを返す"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツ

    # タイトルスライド
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title[:50]
    slide.placeholders[1].text = "AI生成ドラフト"

    # コンテンツをスライドに分割（見出しで分割）
    current_title = "概要"
    current_content = []

    for line in content.split("\n"):
        line = line.strip()
        if line.startswith("## ") or line.startswith("### "):
            # 前のスライドを保存
            if current_content:
                _add_pptx_slide(prs, slide_layout, current_title, current_content)
            current_title = line.lstrip("# ").strip()
            current_content = []
        elif line:
            current_content.append(line)

    # 最後のスライド
    if current_content:
        _add_pptx_slide(prs, slide_layout, current_title, current_content)

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name


def _add_pptx_slide(prs, layout, title: str, content: list[str]):
    """PowerPointにスライドを追加するヘルパー"""
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:50]
    tf = slide.placeholders[1].text_frame
    tf.text = "\n".join(content[:10])  # 最大10行
