"""
ソース管理 APIルーター
ファイル/メモ/パス参照/URLをナレッジソースとして登録・管理し、RAG検索インデックス（documents / documents_fts）に反映する
"""
import os
import re
import json
import html
import sqlite3
import hashlib
import unicodedata
import uuid
from pathlib import Path

import httpx
from fastapi import APIRouter, HTTPException, UploadFile, File
from pydantic import BaseModel
from loguru import logger

from services.llm_client import llm_client

router = APIRouter()

# 実ファイルをSOURCES_DIR配下に保有しており、削除時に物理削除してよいsource_type
MANAGED_SOURCE_TYPES = {"upload", "memo", "url"}
VALID_PATH_SOURCE_TYPES = {"local_path", "server_path", "auto_search"}

FALLBACK_SUGGESTIONS = [
    "登録されているソースの主要なポイントを要約してください。",
    "最近追加されたソースについて詳しく教えてください。",
    "ソースの内容に関するよくある質問は何ですか？",
]

AUTO_SEARCH_LIMIT = 50

DB_PATH = os.getenv("SQLITE_DB_PATH", "/data/sqlite/knowledge.db")
SOURCES_DIR = Path(os.getenv("SOURCES_DIR", "/data/sources"))

ALLOWED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".txt": "txt",
}


def _ensure_storage():
    """sourcesテーブルと保存先ディレクトリを用意する"""
    SOURCES_DIR.mkdir(parents=True, exist_ok=True)
    Path(DB_PATH).parent.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(DB_PATH)
    try:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS sources (
                id           INTEGER PRIMARY KEY AUTOINCREMENT,
                name         TEXT NOT NULL,
                type         TEXT NOT NULL,
                size         INTEGER NOT NULL DEFAULT 0,
                file_path    TEXT NOT NULL UNIQUE,
                selected     INTEGER NOT NULL DEFAULT 0,
                uploaded_at  DATETIME DEFAULT CURRENT_TIMESTAMP
            )
        """)
        # source_type列のマイグレーション（既存DBに対する後方互換追加）
        existing_cols = {row[1] for row in conn.execute("PRAGMA table_info(sources)").fetchall()}
        if "source_type" not in existing_cols:
            conn.execute("ALTER TABLE sources ADD COLUMN source_type TEXT")
            conn.execute("UPDATE sources SET source_type = 'memo' WHERE type = 'memo' AND source_type IS NULL")
            conn.execute("UPDATE sources SET source_type = 'upload' WHERE source_type IS NULL")
        conn.commit()
    finally:
        conn.close()


_ensure_storage()


def _extract_text(path: Path, ext: str) -> str:
    """ファイルからテキストを抽出する（RAGインデクス登録用）"""
    try:
        if ext == ".txt":
            return path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        elif ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        texts.append(row_text)
            return "\n".join(texts)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
    except Exception as e:
        logger.warning(f"テキスト抽出エラー [{path.name}]: {e}")
    return ""


def _index_for_rag(file_path: Path, content: str):
    """既存のdocuments / documents_fts テーブルに登録し、RAG検索対象に加える"""
    if not content.strip():
        return
    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.cursor()
        doc_id = str(file_path)
        checksum = hashlib.md5(content.encode("utf-8")).hexdigest()
        cursor.execute("""
            INSERT INTO documents (file_path, file_name, file_type, file_size, checksum)
            VALUES (?, ?, ?, ?, ?)
            ON CONFLICT(file_path) DO UPDATE SET
                checksum = excluded.checksum,
                file_size = excluded.file_size,
                updated_at = CURRENT_TIMESTAMP
        """, (doc_id, file_path.name, file_path.suffix.lower(), file_path.stat().st_size, checksum))
        cursor.execute("DELETE FROM documents_fts WHERE doc_id = ?", (doc_id,))
        cursor.execute("""
            INSERT INTO documents_fts (doc_id, file_path, file_name, content)
            VALUES (?, ?, ?, ?)
        """, (doc_id, doc_id, file_path.name, content))
        conn.commit()
    finally:
        conn.close()

    # ベクトルインデックスへの登録処理を追加
    try:
        from services.vector_engine import vector_engine
        vector_engine.index_document(doc_id, file_path.name, content)
    except Exception as e:
        logger.warning(f"ベクトルインデックス登録に失敗（FTS5は成功）: {e}")


def _unindex(file_path: Path):
    """documents / documents_fts からエントリを削除する"""
    doc_id = str(file_path)
    conn = sqlite3.connect(DB_PATH)
    try:
        conn.execute("DELETE FROM documents WHERE file_path = ?", (doc_id,))
        conn.execute("DELETE FROM documents_fts WHERE doc_id = ?", (doc_id,))
        conn.commit()
    finally:
        conn.close()

    # ベクトルインデックスからの削除処理を追加
    try:
        from services.vector_engine import vector_engine
        vector_engine.remove_document(doc_id)
    except Exception as e:
        logger.warning(f"ベクトルインデックス削除に失敗: {e}")


def _safe_filename(name: str) -> str:
    """パストラバーサル対策込みの保存用ファイル名を生成する"""
    base = re.sub(r"[/\\]", "_", Path(name).name)
    return f"{uuid.uuid4().hex}_{base}"


class SelectedUpdateRequest(BaseModel):
    selected: bool


class MemoRequest(BaseModel):
    title: str
    content: str


class FromPathRequest(BaseModel):
    path: str
    label: str
    source_type: str  # "local_path" | "server_path" | "auto_search"


class FromUrlRequest(BaseModel):
    url: str
    label: str


@router.get("/")
async def list_sources():
    """登録済みソースの一覧を取得する"""
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        rows = conn.execute(
            "SELECT id, name, type, size, uploaded_at, selected FROM sources ORDER BY uploaded_at DESC"
        ).fetchall()
        return {
            "sources": [
                {
                    "id": r["id"],
                    "name": r["name"],
                    "type": r["type"],
                    "size": r["size"],
                    "uploaded_at": r["uploaded_at"],
                    "selected": bool(r["selected"]),
                }
                for r in rows
            ]
        }
    finally:
        conn.close()


@router.get("/suggestions")
async def get_suggestions(mode: str = "auto"):
    """現在のソース一覧をもとにLLMで推奨プロンプトを3件生成する"""
    conn = sqlite3.connect(DB_PATH)
    try:
        rows = conn.execute(
            "SELECT name, type, file_path FROM sources ORDER BY uploaded_at DESC LIMIT 10"
        ).fetchall()
    finally:
        conn.close()

    if not rows:
        return {"suggestions": FALLBACK_SUGGESTIONS}

    doc_descriptions = []
    for name, type_, file_path in rows:
        if type_ == "memo":
            try:
                content = Path(file_path).read_text(encoding="utf-8", errors="ignore")
                doc_descriptions.append(f"{name}（メモ）: {content[:120]}")
            except OSError:
                doc_descriptions.append(f"{name}（メモ）")
        else:
            doc_descriptions.append(f"{name}（{type_}）")

    doc_list_text = "\n".join(f"- {d}" for d in doc_descriptions)
    prompt = (
        "以下のドキュメント一覧をもとに、ユーザーが質問したくなる具体的な質問を3件、"
        "JSON配列で返してください。余分な説明は不要です。\n"
        '例: ["質問1","質問2","質問3"]\n\n'
        f"ドキュメント:\n{doc_list_text}"
    )

    try:
        raw = await llm_client.generate(prompt)
        match = re.search(r"\[.*\]", raw, re.DOTALL)
        if not match:
            raise ValueError("応答にJSON配列が見つかりません")
        suggestions = json.loads(match.group(0))
        if not isinstance(suggestions, list) or not all(isinstance(s, str) for s in suggestions):
            raise ValueError("不正な形式の応答")
        suggestions = [s for s in suggestions if s.strip()][:3]
        if not suggestions:
            raise ValueError("空の配列")
    except Exception as e:
        logger.warning(f"推奨プロンプト生成に失敗、フォールバックを使用します: {e}")
        suggestions = FALLBACK_SUGGESTIONS

    return {"suggestions": suggestions}


@router.get("/auto-search")
async def auto_search(keyword: str):
    """登録済みのwatch-paths配下を、ファイル名にkeywordを含むもの限定で再帰検索する（候補表示のみ）"""
    keyword = keyword.strip()
    if not keyword:
        return {"results": []}

    conn = sqlite3.connect(DB_PATH)
    try:
        watch_paths = conn.execute("SELECT path, label FROM watch_paths").fetchall()
    except sqlite3.OperationalError:
        # settingsルーターが未ロードでテーブル未作成の場合はフォールバック
        watch_paths = []
    finally:
        conn.close()

    # macOSはファイル名をNFD形式（濁点・半濁点を結合文字として分離）で保存するため、
    # Docker経由でos.walkすると「ガ」が「カ+゛(U+3099)」に分解されて返ってくる。
    # キーワードはNFC形式で入力されるため、両者をNFKCに正規化してから比較する。
    keyword_normalized = unicodedata.normalize("NFKC", keyword).lower()
    results = []
    for watch_path, label in watch_paths:
        if not os.path.isdir(watch_path):
            logger.warning(f"watch-pathが見つかりません（スキップ): {watch_path}")
            continue
        for root, _dirs, files in os.walk(watch_path, followlinks=True):
            for file_name in files:
                file_name_normalized = unicodedata.normalize("NFKC", file_name).lower()
                if keyword_normalized in file_name_normalized:
                    results.append({
                        "path": os.path.join(root, file_name),
                        "file_name": file_name,
                        "watch_path_label": label,
                    })
                    if len(results) >= AUTO_SEARCH_LIMIT:
                        break
            if len(results) >= AUTO_SEARCH_LIMIT:
                break
        if len(results) >= AUTO_SEARCH_LIMIT:
            break

    return {"results": results}


@router.post("/upload")
async def upload_source(file: UploadFile = File(...)):
    """ファイルをアップロードしてソースとして登録し、RAGインデックスに追加する"""
    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"サポートされていないファイル形式です: {ext or '(拡張子なし)'}")

    content_bytes = await file.read()
    dest_path = SOURCES_DIR / _safe_filename(file.filename)
    dest_path.write_bytes(content_bytes)

    text = _extract_text(dest_path, ext)
    _index_for_rag(dest_path, text)

    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.execute(
            "INSERT INTO sources (name, type, size, file_path, selected, source_type) VALUES (?, ?, ?, ?, 0, 'upload')",
            (file.filename, ALLOWED_EXTENSIONS[ext], len(content_bytes), str(dest_path)),
        )
        conn.commit()
        source_id = cursor.lastrowid
        row = conn.execute(
            "SELECT id, name, type, size, uploaded_at FROM sources WHERE id = ?", (source_id,)
        ).fetchone()
    finally:
        conn.close()

    logger.info(f"ソース追加: {file.filename} ({len(content_bytes)} bytes)")
    return {
        "id": row[0],
        "name": row[1],
        "type": row[2],
        "size": row[3],
        "uploaded_at": row[4],
    }


@router.patch("/{source_id}")
async def update_source_selection(source_id: int, request: SelectedUpdateRequest):
    """ソースの選択状態（チャットコンテキストに含めるか）を更新する"""
    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.execute(
            "UPDATE sources SET selected = ? WHERE id = ?",
            (1 if request.selected else 0, source_id),
        )
        conn.commit()
        if cursor.rowcount == 0:
            raise HTTPException(status_code=404, detail=f"ソースが見つかりません: id={source_id}")
    finally:
        conn.close()
    return {"id": source_id, "selected": request.selected}


@router.delete("/{source_id}")
async def delete_source(source_id: int):
    """ソースを削除し、RAGインデックスからも除外する（パス参照ソースの場合は実ファイルは削除しない）"""
    conn = sqlite3.connect(DB_PATH)
    try:
        row = conn.execute(
            "SELECT file_path, source_type FROM sources WHERE id = ?", (source_id,)
        ).fetchone()
        if row is None:
            raise HTTPException(status_code=404, detail=f"ソースが見つかりません: id={source_id}")
        file_path = Path(row[0])
        source_type = row[1]
        conn.execute("DELETE FROM sources WHERE id = ?", (source_id,))
        conn.commit()
    finally:
        conn.close()

    _unindex(file_path)
    if source_type in MANAGED_SOURCE_TYPES and file_path.exists():
        try:
            file_path.unlink()
        except OSError as e:
            logger.warning(f"ソースファイルの削除に失敗しました [{file_path}]: {e}")

    return {"deleted": True}


@router.post("/memo")
async def create_memo_source(request: MemoRequest):
    """メモをテキストソースとして保存し、RAGインデックスに追加する"""
    dest_path = SOURCES_DIR / _safe_filename(f"{request.title}.txt")
    dest_path.write_text(request.content, encoding="utf-8")

    _index_for_rag(dest_path, request.content)

    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.execute(
            "INSERT INTO sources (name, type, size, file_path, selected, source_type) VALUES (?, ?, ?, ?, 0, 'memo')",
            (request.title, "memo", len(request.content.encode("utf-8")), str(dest_path)),
        )
        conn.commit()
        source_id = cursor.lastrowid
        row = conn.execute(
            "SELECT id, name, uploaded_at FROM sources WHERE id = ?", (source_id,)
        ).fetchone()
    finally:
        conn.close()

    logger.info(f"メモソース追加: {request.title}")
    return {
        "id": row[0],
        "name": row[1],
        "type": "memo",
        "uploaded_at": row[2],
    }


@router.post("/from-path")
async def create_source_from_path(request: FromPathRequest):
    """ファイルをコピーせず、パス参照のままソースとして登録する（ローカルPATH/ファイルサーバPATH/自動検索候補から使用）"""
    if request.source_type not in VALID_PATH_SOURCE_TYPES:
        raise HTTPException(
            status_code=400,
            detail=f"source_typeはlocal_path/server_path/auto_searchのいずれかである必要があります: {request.source_type}",
        )

    path = request.path
    # path はファイルシステム上の実パス（macOSではNFD形式）のため正規化せずそのまま使う。
    # 表示名(label)のみNFKC正規化し、auto-search結果から登録した場合の表記を揃える。
    label = unicodedata.normalize("NFKC", request.label)
    ext = Path(path).suffix.lower()
    type_label = ALLOWED_EXTENSIONS.get(ext, ext.lstrip(".") or "file")

    size = 0
    if os.path.isfile(path):
        size = os.path.getsize(path)
        if ext in ALLOWED_EXTENSIONS:
            text = _extract_text(Path(path), ext)
            _index_for_rag(Path(path), text)
        else:
            logger.warning(f"未対応の拡張子のためRAGインデックスはスキップします: {path}")
    else:
        logger.warning(f"パス参照ソース登録時点でファイルが見つかりません（コンテナ未マウントの可能性）: {path}")

    conn = sqlite3.connect(DB_PATH)
    try:
        try:
            cursor = conn.execute(
                "INSERT INTO sources (name, type, size, file_path, selected, source_type) VALUES (?, ?, ?, ?, 0, ?)",
                (label, type_label, size, path, request.source_type),
            )
        except sqlite3.IntegrityError:
            raise HTTPException(status_code=400, detail=f"このパスは既に登録されています: {path}")
        conn.commit()
        source_id = cursor.lastrowid
        row = conn.execute(
            "SELECT id, name, type, uploaded_at FROM sources WHERE id = ?", (source_id,)
        ).fetchone()
    finally:
        conn.close()

    logger.info(f"パス参照ソース追加: {label} ({path}, {request.source_type})")
    return {
        "id": row[0],
        "name": row[1],
        "type": row[2],
        "path": path,
        "uploaded_at": row[3],
    }


def _strip_html(raw_html: str) -> str:
    """HTMLから簡易的にテキストのみを抽出する"""
    text = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", raw_html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = html.unescape(text)
    return re.sub(r"[ \t]+", " ", text).strip()


@router.post("/from-url")
async def create_source_from_url(request: FromUrlRequest):
    """WEBページのURLからコンテンツを取得し、テキスト抽出してソースとして登録する"""
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
            response = await client.get(request.url)
            response.raise_for_status()
    except httpx.HTTPError as e:
        raise HTTPException(status_code=400, detail=f"URLの取得に失敗しました: {e}")

    content_type = response.headers.get("content-type", "")
    raw_text = response.text
    text = _strip_html(raw_text) if "html" in content_type else raw_text.strip()

    if not text:
        raise HTTPException(status_code=400, detail="URLからテキストを抽出できませんでした")

    dest_path = SOURCES_DIR / _safe_filename(f"{request.label}.txt")
    dest_path.write_text(text, encoding="utf-8")
    _index_for_rag(dest_path, text)

    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.execute(
            "INSERT INTO sources (name, type, size, file_path, selected, source_type) VALUES (?, ?, ?, ?, 0, 'url')",
            (request.label, "url", len(text.encode("utf-8")), str(dest_path)),
        )
        conn.commit()
        source_id = cursor.lastrowid
        row = conn.execute(
            "SELECT id, name, uploaded_at FROM sources WHERE id = ?", (source_id,)
        ).fetchone()
    finally:
        conn.close()

    logger.info(f"URLソース追加: {request.label} ({request.url})")
    return {
        "id": row[0],
        "name": row[1],
        "type": "url",
        "uploaded_at": row[2],
    }
