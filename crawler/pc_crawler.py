"""
PC内部ファイル監視クローラー
watchdog でファイル変更を検知し、SQLite FTS5 に自動インデクシングする

対応ファイル形式: .pdf / .docx / .xlsx / .pptx / .txt
"""
import os
import sqlite3
import hashlib
import time
from pathlib import Path
from datetime import datetime

import httpx
import threading
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from loguru import logger
from dotenv import load_dotenv

load_dotenv()

# ============================================
# 設定
# ============================================
WATCH_FOLDERS = os.getenv("WATCH_FOLDERS", "/watch").split(",")
SQLITE_DB_PATH = os.getenv("SQLITE_DB_PATH", "/data/sqlite/knowledge.db")
API_URL = os.getenv("API_URL", "http://api:8000")
POLL_INTERVAL_SECONDS = int(os.getenv("POLL_INTERVAL_SECONDS", "300"))

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}


# ============================================
# テキスト抽出
# ============================================

def extract_text(file_path: Path) -> str:
    """ファイルからテキストを抽出する"""
    ext = file_path.suffix.lower()
    try:
        if ext == ".txt" or ext == ".md":
            return file_path.read_text(encoding="utf-8", errors="ignore")

        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            texts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
            return "\n".join(texts)

        elif ext == ".docx":
            from docx import Document
            doc = Document(str(file_path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        texts.append(row_text)
            return "\n".join(texts)

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)

    except Exception as e:
        logger.warning(f"テキスト抽出エラー [{file_path.name}]: {e}")
        return ""

    return ""


def get_checksum(file_path: Path) -> str:
    """ファイルのMD5チェックサムを計算する"""
    h = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


# ============================================
# SQLite インデクシング
# ============================================

def index_file(file_path: Path):
    """ファイルをSQLite FTS5にインデクシングする"""
    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return

    logger.info(f"インデクシング: {file_path.name}")

    try:
        checksum = get_checksum(file_path)
        content = extract_text(file_path)
        if not content.strip():
            logger.warning(f"テキスト抽出結果が空: {file_path.name}")
            return

        conn = sqlite3.connect(SQLITE_DB_PATH)
        try:
            cursor = conn.cursor()
            doc_id = str(file_path)

            # 既存レコードを確認
            cursor.execute(
                "SELECT checksum FROM documents WHERE file_path = ?",
                (doc_id,)
            )
            row = cursor.fetchone()

            if row and row[0] == checksum:
                logger.debug(f"変更なし（スキップ）: {file_path.name}")
                return

            # documentsテーブルを更新
            cursor.execute("""
                INSERT INTO documents (file_path, file_name, file_type, file_size, checksum)
                VALUES (?, ?, ?, ?, ?)
                ON CONFLICT(file_path) DO UPDATE SET
                    checksum = excluded.checksum,
                    file_size = excluded.file_size,
                    updated_at = CURRENT_TIMESTAMP
            """, (
                doc_id,
                file_path.name,
                ext,
                file_path.stat().st_size,
                checksum,
            ))

            # FTS5インデクスを更新
            cursor.execute(
                "DELETE FROM documents_fts WHERE doc_id = ?",
                (doc_id,)
            )
            cursor.execute("""
                INSERT INTO documents_fts (doc_id, file_path, file_name, content)
                VALUES (?, ?, ?, ?)
            """, (doc_id, str(file_path), file_path.name, content))

            conn.commit()
            logger.success(f"✅ インデクシング完了: {file_path.name} ({len(content)}文字)")

        finally:
            conn.close()

    except Exception as e:
        logger.error(f"インデクシングエラー [{file_path.name}]: {e}")


def remove_from_index(file_path: Path):
    """ファイルをインデクスから削除する"""
    doc_id = str(file_path)
    conn = sqlite3.connect(SQLITE_DB_PATH)
    try:
        cursor = conn.cursor()
        cursor.execute("DELETE FROM documents WHERE file_path = ?", (doc_id,))
        cursor.execute("DELETE FROM documents_fts WHERE doc_id = ?", (doc_id,))
        conn.commit()
        logger.info(f"インデクスから削除: {file_path.name}")
    finally:
        conn.close()


# ============================================
# ファイル監視イベントハンドラ
# ============================================

class DocumentEventHandler(FileSystemEventHandler):
    """watchdog イベントハンドラ"""

    def on_created(self, event):
        if not event.is_directory:
            index_file(Path(event.src_path))

    def on_modified(self, event):
        if not event.is_directory:
            index_file(Path(event.src_path))

    def on_deleted(self, event):
        if not event.is_directory:
            remove_from_index(Path(event.src_path))

    def on_moved(self, event):
        if not event.is_directory:
            remove_from_index(Path(event.src_path))
            index_file(Path(event.dest_path))


# ============================================
# メイン
# ============================================

def initial_scan():
    """起動時および定期実行時に監視フォルダ内の既存ファイルをすべてスキャンする"""
    logger.info("🔍 スキャン開始（差分検知）...")
    count = 0
    
    # 1. 既存ファイルをインデクシング
    for folder in WATCH_FOLDERS:
        folder = folder.strip()
        if not Path(folder).exists():
            logger.warning(f"監視フォルダが存在しません: {folder}")
            continue
        for file_path in Path(folder).rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                index_file(file_path)
                count += 1
                
    # 2. 実ファイルが存在しない古いインデックスを削除するクリーンアップ
    cleanup_count = 0
    try:
        conn = sqlite3.connect(SQLITE_DB_PATH)
        cursor = conn.cursor()
        cursor.execute("SELECT file_path FROM documents")
        db_files = [row[0] for row in cursor.fetchall()]
        conn.close()
        
        for db_file in db_files:
            if not Path(db_file).exists():
                remove_from_index(Path(db_file))
                cleanup_count += 1
    except Exception as e:
        logger.error(f"クリーンアップスキャンエラー: {e}")
        
    logger.success(f"✅ スキャン完了: 新規/更新 {count}件, 削除クリーンアップ {cleanup_count}件")


def polling_scan_loop():
    """定期的に監視フォルダ内をポーリングスキャンし、差分を反映する（ファイルサーバー等用）"""
    logger.info(f"⏳ 定期ポーリングスキャンが有効化されました (間隔: {POLL_INTERVAL_SECONDS}秒)")
    while True:
        try:
            time.sleep(POLL_INTERVAL_SECONDS)
            logger.info("🔍 定期ポーリングスキャンを開始します...")
            initial_scan()
        except Exception as e:
            logger.error(f"定期ポーリングスキャン中にエラーが発生しました: {e}")


def main():
    logger.info("🚀 PC内部・ファイルサーバークローラー起動")
    logger.info(f"監視フォルダ: {WATCH_FOLDERS}")
    logger.info(f"対応形式: {SUPPORTED_EXTENSIONS}")
    logger.info(f"ポーリング間隔: {POLL_INTERVAL_SECONDS}秒")
 
    # DB初期化待機（APIコンテナが先に起動するため）
    time.sleep(5)
 
    # 初期スキャン
    initial_scan()
 
    # 定期ポーリングスキャンを別スレッドで開始
    if POLL_INTERVAL_SECONDS > 0:
        polling_thread = threading.Thread(target=polling_scan_loop, daemon=True)
        polling_thread.start()
 
    # watchdog でリアルタイム監視
    event_handler = DocumentEventHandler()
    observer = Observer()
 
    for folder in WATCH_FOLDERS:
        folder = folder.strip()
        if Path(folder).exists():
            observer.schedule(event_handler, folder, recursive=True)
            logger.info(f"👁️  監視開始: {folder}")
 
    observer.start()
    logger.info("✅ ファイル監視中... (Ctrl+C で停止)")
 
    try:
        while True:
            time.sleep(10)
    except KeyboardInterrupt:
        observer.stop()
        logger.info("🛑 クローラー停止")
 
    observer.join()


if __name__ == "__main__":
    main()
